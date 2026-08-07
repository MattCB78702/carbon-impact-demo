"""Parse CarbonBetter's New Belgium catalogue deck into structured project records.

READ-ONLY on the source PPTX. Writes only extract/catalogue_projects.json.

The deck uses inconsistent field labels across slides ("Project ID:", "Registry:",
"Registry (ID):", "Certifications:", "Ratings/ Certifications:"), so labels are
normalized to a single schema here rather than downstream.
"""
import json
import os
import re

from pptx import Presentation

SRC = (r"C:\Users\MattHendren\OneDrive - CarbonBetter\Claude\Carbon Credit Product"
       r"\New Belgium\New Belgium Brewing Catalogue Carbon offset_July 2026.pptx")
OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# Project slides only. 1-5 are cover/about/approach/clients/criteria; 32 is next steps.
FIRST_PROJECT_SLIDE = 6
LAST_PROJECT_SLIDE = 31

LABEL_MAP = {
    "project id": "registry_id", "registry": "registry_id", "registry (id)": "registry_id",
    "vintage": "vintage",
    "credit type": "credit_type",
    "location": "location",
    "tech": "tech",
    "verified by": "verifier",
    "methodology": "methodology",
    "certifications": "ratings", "ratings": "ratings",
    "ratings/ certifications": "ratings", "ratings / certifications": "ratings",
    "volume": "volume",
    "price ($/ton)": "price", "price": "price",
}

COUNTRY_HINTS = ["USA", "Canada", "Mexico"]


def norm_label(raw):
    key = raw.strip().rstrip(":").strip().lower()
    key = re.sub(r"\s+", " ", key)
    return LABEL_MAP.get(key)


def clean(text):
    return re.sub(r"\s+", " ", (text or "").replace("\xa0", " ")).strip()


def main():
    prs = Presentation(SRC)
    projects = []

    for idx, slide in enumerate(prs.slides, start=1):
        if not (FIRST_PROJECT_SLIDE <= idx <= LAST_PROJECT_SLIDE):
            continue

        fields, prose, country = {}, [], ""

        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [clean(c.text) for c in row.cells]
                    if len(cells) < 2:
                        continue
                    key = norm_label(cells[0])
                    if key and cells[1]:
                        # keep the first non-empty value if a label repeats
                        fields.setdefault(key, cells[1])
            elif shape.has_text_frame:
                t = clean(shape.text_frame.text)
                if not t:
                    continue
                if t in COUNTRY_HINTS:
                    country = t
                elif re.fullmatch(r"(Also includes )?Removal credits", t, re.I):
                    fields["removal_flag"] = t
                else:
                    prose.append(t)

        if not fields and not prose:
            continue

        # Longest prose block is the description; the shortest leading one is the title.
        prose_sorted = sorted(prose, key=len, reverse=True)
        description = prose_sorted[0] if prose_sorted else ""
        title_candidates = [p for p in prose if p != description and len(p) < 120]
        name = title_candidates[0] if title_candidates else f"Slide {idx}"

        projects.append({
            "slide": idx,
            "name": name,
            "country": country,
            "registry_id": fields.get("registry_id", ""),
            "vintage": fields.get("vintage", ""),
            "credit_type": fields.get("credit_type", ""),
            "location": fields.get("location", ""),
            "tech": fields.get("tech", ""),
            "methodology": fields.get("methodology", ""),
            "verifier": fields.get("verifier", ""),
            "ratings": fields.get("ratings", ""),
            "volume": fields.get("volume", ""),
            "price": fields.get("price", ""),
            "removal_flag": fields.get("removal_flag", ""),
            "description": description,
        })

    with open(os.path.join(OUT_DIR, "catalogue_projects.json"), "w", encoding="utf-8") as f:
        json.dump(projects, f, indent=2, ensure_ascii=False)

    print(f"slides in deck: {len(prs.slides)}")
    print(f"project slides parsed: {len(projects)}")
    missing = [p["slide"] for p in projects if not p["registry_id"] or not p["price"]]
    print(f"slides missing registry_id or price: {missing or 'none'}")
    for p in projects:
        print(f"  s{p['slide']:>2} | {p['name'][:52]:<52} | {p['registry_id'][:34]:<34} | {p['tech'][:26]:<26} | {p['price'][:22]}")


if __name__ == "__main__":
    main()
